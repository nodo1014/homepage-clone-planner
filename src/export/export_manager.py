"""
내보내기 관리자 모듈

이 모듈은 클론 기획서를 다양한 형식(Markdown, PPT, ZIP 등)으로 내보내는 기능을 제공합니다.
"""
import os
import shutil
import logging
import zipfile
import json
import markdown
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from typing import Dict, Any, List, Optional, Union, Tuple
import tempfile
import re
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse, urljoin
from weasyprint import HTML  # PDF 변환을 위한 라이브러리 추가

# 내부 모듈 임포트
from src.api.notion_client import NotionClient
from src.api.gdrive_client import GoogleDriveClient  # 구글 드라이브 클라이언트 임포트

# 로거 설정
logger = logging.getLogger(__name__)

class ExportManager:
    """다양한 형식으로 내보내기를 관리하는 클래스"""
    
    def __init__(self, output_dir: Optional[Union[str, Path]] = None):
        """
        내보내기 관리자 초기화
        
        Args:
            output_dir (Optional[Union[str, Path]]): 내보내기 결과물이 저장될 디렉토리 (None이면 임시 디렉토리 사용)
        """
        # 로거 설정
        self.logger = logger
        
        if output_dir:
            self.output_dir = Path(output_dir)
            self.output_dir.mkdir(parents=True, exist_ok=True)
        else:
            self.temp_dir = tempfile.TemporaryDirectory()
            self.output_dir = Path(self.temp_dir.name)
            
        logger.info(f"내보내기 디렉토리 설정: {self.output_dir}")
        
        # 지원하는 내보내기 형식
        self.supported_formats = ['md', 'html', 'pptx', 'zip', 'json', 'ai_meta', 'pdf', 'notion', 'gdrive']  # gdrive 추가
    
    def __del__(self):
        """임시 디렉토리 정리"""
        if hasattr(self, 'temp_dir'):
            self.temp_dir.cleanup()
    
    def export_to_format(self, content: Dict[str, Any], format_type: str, 
                         filename: str = "export", **kwargs) -> str:
        """
        지정된 형식으로 내용 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠 (기획서 데이터)
            format_type (str): 내보내기 형식 ('md', 'html', 'pptx', 'zip', 'json', 'ai_meta', 'pdf', 'notion', 'gdrive')
            filename (str): 출력 파일 이름 (확장자 제외)
            **kwargs: 형식별 추가 옵션
            
        Returns:
            str: 생성된 파일의 경로 또는 URL
        """
        format_type = format_type.lower()
        
        if format_type not in self.supported_formats:
            raise ValueError(f"지원하지 않는 형식: {format_type}. 지원 형식: {', '.join(self.supported_formats)}")
        
        # 형식에 따라 적절한 내보내기 방법 호출
        if format_type == 'md':
            return self.export_to_markdown(content, filename, **kwargs)
        elif format_type == 'html':
            return self.export_to_html(content, filename, **kwargs)
        elif format_type == 'pptx':
            return self.export_to_pptx(content, filename, **kwargs)
        elif format_type == 'zip':
            return self.export_to_zip(content, filename, **kwargs)
        elif format_type == 'json':
            return self.export_to_json(content, filename, **kwargs)
        elif format_type == 'ai_meta':
            return self.export_for_ai_analysis(content, filename=filename, **kwargs)
        elif format_type == 'pdf':
            return self.export_to_pdf(content, filename, **kwargs)
        elif format_type == 'notion':
            return self.export_to_notion(content, **kwargs)
        elif format_type == 'gdrive':
            return self.export_to_google_drive(content, filename, **kwargs)
    
    def export_to_markdown(self, content: Dict[str, Any], filename: str = "export", 
                           template: Optional[str] = None, **kwargs) -> str:
        """
        마크다운 형식으로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            template (Optional[str]): 마크다운 템플릿 경로 (None이면 기본 템플릿 사용)
            
        Returns:
            str: 생성된 마크다운 파일 경로
        """
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}.md"
        
        try:
            # 기본 템플릿 또는 사용자 지정 템플릿 로드
            md_content = self._generate_markdown_content(content, template)
            
            # 파일 작성
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(md_content)
            
            logger.info(f"마크다운 파일 생성: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"마크다운 내보내기 실패: {str(e)}")
            raise
    
    def export_to_html(self, content: Dict[str, Any], filename: str = "export", 
                       css: Optional[str] = None, **kwargs) -> str:
        """
        HTML 형식으로 내보내기 (마크다운을 HTML로 변환)
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            css (Optional[str]): CSS 스타일시트 경로
            
        Returns:
            str: 생성된 HTML 파일 경로
        """
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}.html"
        
        try:
            # 먼저 마크다운 생성
            md_content = self._generate_markdown_content(content)
            
            # 마크다운을 HTML로 변환
            html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
            
            # CSS 추가
            if css and os.path.exists(css):
                with open(css, "r", encoding="utf-8") as f:
                    css_content = f.read()
                html_template = f"""<!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <meta name="viewport" content="width=device-width, initial-scale=1.0">
                    <title>{content.get('title', '클론 기획서')}</title>
                    <style>
                    {css_content}
                    </style>
                </head>
                <body>
                    {html_content}
                </body>
                </html>
                """
            else:
                # 기본 스타일
                html_template = f"""<!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <meta name="viewport" content="width=device-width, initial-scale=1.0">
                    <title>{content.get('title', '클론 기획서')}</title>
                    <style>
                    body {{
                        font-family: 'Noto Sans KR', Arial, sans-serif;
                        line-height: 1.6;
                        max-width: 900px;
                        margin: 0 auto;
                        padding: 20px;
                        color: #333;
                    }}
                    h1, h2, h3, h4, h5, h6 {{
                        margin-top: 1.5em;
                        margin-bottom: 0.5em;
                        color: #1a1a1a;
                    }}
                    h1 {{ font-size: 2.2em; border-bottom: 1px solid #eee; padding-bottom: 0.3em; }}
                    h2 {{ font-size: 1.8em; border-bottom: 1px solid #eee; padding-bottom: 0.3em; }}
                    h3 {{ font-size: 1.5em; }}
                    h4 {{ font-size: 1.3em; }}
                    table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
                    th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
                    th {{ background-color: #f2f2f2; }}
                    img {{ max-width: 100%; height: auto; }}
                    code {{ background-color: #f5f5f5; padding: 2px 4px; border-radius: 3px; }}
                    pre {{ background-color: #f5f5f5; padding: 1em; overflow-x: auto; border-radius: 3px; }}
                    blockquote {{ background-color: #f9f9f9; border-left: 4px solid #ccc; margin: 1em 0; padding: 0.5em 1em; }}
                    </style>
                </head>
                <body>
                    {html_content}
                </body>
                </html>
                """
            
            # 파일 작성
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(html_template)
            
            logger.info(f"HTML 파일 생성: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"HTML 내보내기 실패: {str(e)}")
            raise
    
    def export_to_pptx(self, content: Dict[str, Any], filename: str = "export", 
                       template_pptx: Optional[str] = None, **kwargs) -> str:
        """
        PowerPoint 형식으로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            template_pptx (Optional[str]): 템플릿 PPTX 파일 경로
            
        Returns:
            str: 생성된 PPTX 파일 경로
        """
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}.pptx"
        
        try:
            # 템플릿 또는 빈 프레젠테이션 로드
            if template_pptx and os.path.exists(template_pptx):
                prs = Presentation(template_pptx)
            else:
                prs = Presentation()
            
            # 프레젠테이션 생성
            self._generate_presentation(prs, content)
            
            # 파일 저장
            prs.save(output_path)
            
            logger.info(f"PowerPoint 파일 생성: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"PowerPoint 내보내기 실패: {str(e)}")
            raise
    
    def export_to_json(self, content: Dict[str, Any], filename: str = "export", 
                       indent: int = 2, **kwargs) -> str:
        """
        JSON 형식으로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            indent (int): JSON 들여쓰기 수준
            
        Returns:
            str: 생성된 JSON 파일 경로
        """
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}.json"
        
        try:
            # JSON 형식으로 저장
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(content, f, ensure_ascii=False, indent=indent)
            
            logger.info(f"JSON 파일 생성: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"JSON 내보내기 실패: {str(e)}")
            raise
    
    def export_to_zip(self, content: Dict[str, Any], filename: str = "export", 
                      include_formats: List[str] = None, **kwargs) -> str:
        """
        여러 형식을 ZIP 아카이브로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            include_formats (List[str]): ZIP에 포함할 형식 목록 (기본: md, html, pptx, json, pdf)
            
        Returns:
            str: 생성된 ZIP 파일 경로
        """
        # 기본 포함 형식
        if include_formats is None:
            include_formats = ['md', 'html', 'pptx', 'json', 'pdf']  # notion은 URL만 반환하므로 ZIP에 포함하지 않음
        
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}.zip"
        
        # 내보내기를 위한 임시 디렉토리 생성
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            try:
                # 각 형식으로 내보내기 수행
                exported_files = []
                for fmt in include_formats:
                    if fmt in self.supported_formats and fmt != 'zip':
                        export_path = self.export_to_format(
                            content=content,
                            format_type=fmt,
                            filename=filename,
                            output_dir=temp_path,
                            **kwargs
                        )
                        exported_files.append(Path(export_path))
                
                # 이미지 및 기타 자원 복사 (있는 경우)
                if 'resources' in content and isinstance(content['resources'], dict):
                    resources_dir = temp_path / "resources"
                    resources_dir.mkdir(exist_ok=True)
                    
                    for res_name, res_path in content['resources'].items():
                        if os.path.exists(res_path):
                            dest_path = resources_dir / os.path.basename(res_path)
                            shutil.copy2(res_path, dest_path)
                            exported_files.append(dest_path)
                
                # ZIP 파일 생성
                with zipfile.ZipFile(output_path, 'w') as zip_file:
                    for file_path in exported_files:
                        zip_file.write(
                            file_path, 
                            arcname=file_path.name
                        )
                
                logger.info(f"ZIP 아카이브 생성: {output_path}")
                return str(output_path)
                
            except Exception as e:
                logger.error(f"ZIP 내보내기 실패: {str(e)}")
                raise
    
    def export_for_ai_analysis(self, content: Dict[str, Any], url: str = None, 
                              filename: str = "ai_analysis", **kwargs) -> str:
        """
        AI 분석을 위한 형식으로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            url (str): 분석할 웹사이트 URL (None이면 content에서 추출)
            filename (str): 출력 파일 이름 (확장자 제외)
            
        Returns:
            str: 생성된 AI 메타데이터 파일 경로
        """
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}_ai_meta.json"
        
        try:
            # URL이 제공되지 않은 경우 content에서 추출
            if not url and 'website' in content and 'url' in content['website']:
                url = content['website']['url']
            
            if not url:
                raise ValueError("AI 분석을 위한 URL이 제공되지 않았습니다.")
            
            # 웹사이트 분석 메타데이터 생성
            ai_metadata = self._generate_ai_metadata(content, url)
            
            # JSON 형식으로 저장
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(ai_metadata, f, ensure_ascii=False, indent=2)
            
            logger.info(f"AI 분석용 메타데이터 파일 생성: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"AI 분석용 메타데이터 내보내기 실패: {str(e)}")
            raise
    
    def _generate_ai_metadata(self, content: Dict[str, Any], url: str) -> Dict[str, Any]:
        """
        AI 분석을 위한 메타데이터 생성
        
        Args:
            content (Dict[str, Any]): 기획서 콘텐츠
            url (str): 분석할 웹사이트 URL
            
        Returns:
            Dict[str, Any]: AI 분석용 메타데이터
        """
        ai_metadata = {
            "version": "1.0",
            "generation_date": self._get_current_timestamp(),
            "url": url,
            "title": content.get('title', '클론 기획서'),
            "website_info": content.get('website', {}),
            "ai_analysis_hints": {
                "focus_areas": [
                    "디자인 일관성",
                    "사용자 경험",
                    "성능 최적화",
                    "접근성",
                    "모바일 응답성"
                ],
                "expected_improvements": content.get('development_recommendations', '')
            }
        }
        
        # 웹사이트 분석 시도 (가능한 경우)
        try:
            site_structure = self._analyze_website_structure(url)
            ai_metadata["site_structure"] = site_structure
        except Exception as e:
            logger.warning(f"웹사이트 구조 분석 실패: {str(e)}")
            ai_metadata["site_structure"] = {"error": str(e)}
        
        # 페이지 구조 정보 추가
        if 'page_structure' in content and isinstance(content['page_structure'], list):
            ai_metadata["page_structure"] = content['page_structure']
        
        # 디자인 분석 정보 추가
        if 'design_analysis' in content and isinstance(content['design_analysis'], dict):
            ai_metadata["design_analysis"] = content['design_analysis']
        
        # 기능 분석 정보 추가
        if 'functional_analysis' in content and isinstance(content['functional_analysis'], dict):
            ai_metadata["functional_analysis"] = content['functional_analysis']
        
        return ai_metadata
    
    def _analyze_website_structure(self, url: str, max_depth: int = 2) -> Dict[str, Any]:
        """
        웹사이트의 구조를 분석
        
        Args:
            url (str): 분석할 웹사이트 URL
            max_depth (int): 분석할 최대 깊이
            
        Returns:
            Dict[str, Any]: 웹사이트 구조 정보
        """
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }
            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # 기본 메타데이터 추출
            title = soup.title.text.strip() if soup.title else "제목 없음"
            meta_description = ""
            if soup.find("meta", attrs={"name": "description"}):
                meta_description = soup.find("meta", attrs={"name": "description"}).get("content", "")
            
            # 주요 구조 요소 식별
            structure = {
                "url": url,
                "title": title,
                "meta_description": meta_description,
                "headers": self._extract_headers(soup),
                "navigation": self._extract_navigation(soup),
                "main_sections": self._extract_main_sections(soup),
                "forms": self._extract_forms(soup),
                "images": self._extract_image_stats(soup),
                "links": self._extract_links(soup, url, max_depth)
            }
            
            return structure
            
        except Exception as e:
            logger.error(f"웹사이트 구조 분석 중 오류 발생: {str(e)}")
            return {"error": str(e)}
    
    def _extract_headers(self, soup: BeautifulSoup) -> List[Dict[str, str]]:
        """
        페이지에서 헤더 요소 추출
        
        Args:
            soup (BeautifulSoup): 파싱된 HTML
            
        Returns:
            List[Dict[str, str]]: 헤더 정보 목록
        """
        headers = []
        for i in range(1, 7):
            for header in soup.find_all(f'h{i}'):
                headers.append({
                    "level": i,
                    "text": header.get_text(strip=True)
                })
        return headers
    
    def _extract_navigation(self, soup: BeautifulSoup) -> List[Dict[str, str]]:
        """
        페이지에서 내비게이션 요소 추출
        
        Args:
            soup (BeautifulSoup): 파싱된 HTML
            
        Returns:
            List[Dict[str, str]]: 내비게이션 정보
        """
        nav_items = []
        
        # 일반적인 네비게이션 요소 찾기
        nav_elements = soup.find_all(['nav', 'div'], class_=lambda x: x and ('nav' in x.lower() or 'menu' in x.lower()))
        
        for nav in nav_elements:
            links = nav.find_all('a')
            for link in links:
                nav_items.append({
                    "text": link.get_text(strip=True),
                    "href": link.get('href', '#')
                })
        
        return nav_items
    
    def _extract_main_sections(self, soup: BeautifulSoup) -> List[Dict[str, Any]]:
        """
        페이지의 주요 섹션 추출
        
        Args:
            soup (BeautifulSoup): 파싱된 HTML
            
        Returns:
            List[Dict[str, Any]]: 주요 섹션 정보
        """
        sections = []
        
        # 시맨틱 요소 확인
        semantic_elements = soup.find_all(['section', 'article', 'main', 'aside', 'header', 'footer'])
        
        for element in semantic_elements:
            section_info = {
                "type": element.name,
                "id": element.get('id', ''),
                "classes": element.get('class', []),
                "text_length": len(element.get_text(strip=True)),
                "has_images": len(element.find_all('img')) > 0,
                "has_links": len(element.find_all('a')) > 0
            }
            
            # 제목 요소가 있는지 확인
            heading = element.find(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
            if heading:
                section_info["heading"] = heading.get_text(strip=True)
            
            sections.append(section_info)
        
        return sections
    
    def _extract_forms(self, soup: BeautifulSoup) -> List[Dict[str, Any]]:
        """
        페이지의 폼 요소 추출
        
        Args:
            soup (BeautifulSoup): 파싱된 HTML
            
        Returns:
            List[Dict[str, Any]]: 폼 정보
        """
        forms = []
        
        for form in soup.find_all('form'):
            form_info = {
                "action": form.get('action', ''),
                "method": form.get('method', 'get'),
                "id": form.get('id', ''),
                "fields": []
            }
            
            for input_element in form.find_all(['input', 'textarea', 'select']):
                field_info = {
                    "type": input_element.name,
                    "name": input_element.get('name', ''),
                    "id": input_element.get('id', '')
                }
                
                if input_element.name == 'input':
                    field_info["input_type"] = input_element.get('type', 'text')
                
                form_info["fields"].append(field_info)
            
            forms.append(form_info)
        
        return forms
    
    def _extract_image_stats(self, soup: BeautifulSoup) -> Dict[str, Any]:
        """
        페이지의 이미지 통계 추출
        
        Args:
            soup (BeautifulSoup): 파싱된 HTML
            
        Returns:
            Dict[str, Any]: 이미지 통계
        """
        images = soup.find_all('img')
        
        stats = {
            "count": len(images),
            "with_alt": sum(1 for img in images if img.get('alt')),
            "without_alt": sum(1 for img in images if not img.get('alt'))
        }
        
        return stats
    
    def _extract_links(self, soup: BeautifulSoup, base_url: str, max_depth: int) -> Dict[str, Any]:
        """
        페이지의 링크 정보 추출
        
        Args:
            soup (BeautifulSoup): 파싱된 HTML
            base_url (str): 기본 URL
            max_depth (int): 최대 탐색 깊이
            
        Returns:
            Dict[str, Any]: 링크 정보
        """
        links = soup.find_all('a')
        
        internal_links = []
        external_links = []
        
        base_domain = urlparse(base_url).netloc
        
        for link in links:
            href = link.get('href', '')
            if not href or href.startswith('#'):
                continue
                
            # 상대 URL을 절대 URL로 변환
            absolute_url = urljoin(base_url, href)
            link_domain = urlparse(absolute_url).netloc
            
            link_info = {
                "url": absolute_url,
                "text": link.get_text(strip=True),
                "title": link.get('title', '')
            }
            
            if link_domain == base_domain:
                internal_links.append(link_info)
            else:
                external_links.append(link_info)
        
        return {
            "internal_count": len(internal_links),
            "external_count": len(external_links),
            "internal_sample": internal_links[:10],  # 처음 10개만 샘플로 포함
            "external_sample": external_links[:5]    # 처음 5개만 샘플로 포함
        }
    
    def _get_current_timestamp(self) -> str:
        """현재 타임스탬프를 ISO 형식으로 반환"""
        from datetime import datetime
        return datetime.now().isoformat()
    
    def _generate_markdown_content(self, content: Dict[str, Any], 
                                  template: Optional[str] = None) -> str:
        """
        콘텐츠에서 마크다운 내용 생성
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            template (Optional[str]): 마크다운 템플릿 경로
            
        Returns:
            str: 생성된 마크다운 내용
        """
        # 템플릿이 제공된 경우 로드
        if template and os.path.exists(template):
            with open(template, "r", encoding="utf-8") as f:
                template_content = f.read()
                
            # 템플릿에 변수 채우기
            for key, value in content.items():
                placeholder = f"{{{{{key}}}}}"
                if isinstance(value, str):
                    template_content = template_content.replace(placeholder, value)
                elif isinstance(value, dict):
                    # 중첩된 콘텐츠 처리
                    for nested_key, nested_value in value.items():
                        nested_placeholder = f"{{{{{key}.{nested_key}}}}}"
                        if isinstance(nested_value, str):
                            template_content = template_content.replace(nested_placeholder, nested_value)
            
            return template_content
        
        # 템플릿이 없는 경우 데이터에서 마크다운 생성
        md_parts = []
        
        # 제목
        if 'title' in content:
            md_parts.append(f"# {content['title']}\n")
        
        # 웹사이트 정보
        if 'website' in content:
            website = content['website']
            if isinstance(website, dict):
                md_parts.append("## 웹사이트 정보\n")
                if 'url' in website:
                    md_parts.append(f"- **URL**: {website['url']}")
                if 'name' in website:
                    md_parts.append(f"- **이름**: {website['name']}")
                if 'description' in website:
                    md_parts.append(f"- **설명**: {website['description']}")
                md_parts.append("\n")
        
        # 개요
        if 'overview' in content:
            md_parts.append("## 개요\n")
            md_parts.append(f"{content['overview']}\n")
        
        # 디자인 분석
        if 'design_analysis' in content and isinstance(content['design_analysis'], dict):
            design = content['design_analysis']
            md_parts.append("## 디자인 분석\n")
            
            if 'color_palette' in design:
                md_parts.append("### 색상 팔레트\n")
                if isinstance(design['color_palette'], list):
                    for color in design['color_palette']:
                        md_parts.append(f"- `{color}`")
                else:
                    md_parts.append(design['color_palette'])
                md_parts.append("\n")
            
            if 'typography' in design:
                md_parts.append("### 타이포그래피\n")
                md_parts.append(design['typography'])
                md_parts.append("\n")
            
            if 'layout' in design:
                md_parts.append("### 레이아웃\n")
                md_parts.append(design['layout'])
                md_parts.append("\n")
        
        # 기능 분석
        if 'functional_analysis' in content and isinstance(content['functional_analysis'], dict):
            func = content['functional_analysis']
            md_parts.append("## 기능 분석\n")
            
            if 'key_features' in func and isinstance(func['key_features'], list):
                md_parts.append("### 주요 기능\n")
                for feature in func['key_features']:
                    md_parts.append(f"- {feature}")
                md_parts.append("\n")
            
            if 'user_interactions' in func:
                md_parts.append("### 사용자 인터랙션\n")
                md_parts.append(func['user_interactions'])
                md_parts.append("\n")
        
        # 페이지 구조
        if 'page_structure' in content and isinstance(content['page_structure'], list):
            md_parts.append("## 페이지 구조\n")
            for page in content['page_structure']:
                if isinstance(page, dict) and 'name' in page:
                    md_parts.append(f"### {page['name']}\n")
                    if 'description' in page:
                        md_parts.append(f"{page['description']}\n")
                    if 'components' in page and isinstance(page['components'], list):
                        md_parts.append("#### 구성 요소\n")
                        for component in page['components']:
                            md_parts.append(f"- {component}")
                        md_parts.append("\n")
                else:
                    md_parts.append(f"- {page}\n")
        
        # 기술 스택
        if 'tech_stack' in content and isinstance(content['tech_stack'], list):
            md_parts.append("## 기술 스택\n")
            for tech in content['tech_stack']:
                md_parts.append(f"- {tech}")
            md_parts.append("\n")
        
        # 개발 제안
        if 'development_recommendations' in content:
            md_parts.append("## 개발 제안\n")
            md_parts.append(content['development_recommendations'])
            md_parts.append("\n")
        
        # 결론
        if 'conclusion' in content:
            md_parts.append("## 결론\n")
            md_parts.append(content['conclusion'])
            md_parts.append("\n")
        
        return "\n".join(md_parts)
    
    def _generate_presentation(self, prs: Presentation, content: Dict[str, Any]) -> None:
        """
        콘텐츠를 PowerPoint 프레젠테이션으로 변환
        
        Args:
            prs (Presentation): PowerPoint 프레젠테이션 객체
            content (Dict[str, Any]): 내보낼 콘텐츠
        """
        # 제목 슬라이드
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content.get('title', '클론 기획서')
        subtitle.text = f"웹사이트: {content.get('website', {}).get('name', '')}"
        
        # 개요 슬라이드
        if 'overview' in content:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "개요"
            content_placeholder.text = content['overview']
        
        # 디자인 분석 슬라이드
        if 'design_analysis' in content and isinstance(content['design_analysis'], dict):
            design = content['design_analysis']
            
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "디자인 분석"
            
            text_parts = []
            if 'color_palette' in design:
                text_parts.append("색상 팔레트:")
                if isinstance(design['color_palette'], list):
                    for color in design['color_palette']:
                        text_parts.append(f"• {color}")
                else:
                    text_parts.append(str(design['color_palette']))
            
            if 'typography' in design:
                text_parts.append("\n타이포그래피:")
                text_parts.append(design['typography'])
            
            if 'layout' in design:
                text_parts.append("\n레이아웃:")
                text_parts.append(design['layout'])
            
            content_placeholder.text = "\n".join(text_parts)
        
        # 기능 분석 슬라이드
        if 'functional_analysis' in content and isinstance(content['functional_analysis'], dict):
            func = content['functional_analysis']
            
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "기능 분석"
            
            text_parts = []
            if 'key_features' in func and isinstance(func['key_features'], list):
                text_parts.append("주요 기능:")
                for feature in func['key_features']:
                    text_parts.append(f"• {feature}")
            
            if 'user_interactions' in func:
                text_parts.append("\n사용자 인터랙션:")
                text_parts.append(func['user_interactions'])
            
            content_placeholder.text = "\n".join(text_parts)
        
        # 페이지 구조 슬라이드
        if 'page_structure' in content and isinstance(content['page_structure'], list):
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "페이지 구조"
            
            text_parts = []
            for page in content['page_structure']:
                if isinstance(page, dict) and 'name' in page:
                    text_parts.append(f"• {page['name']}")
                    if 'components' in page and isinstance(page['components'], list):
                        for component in page['components']:
                            text_parts.append(f"  - {component}")
                else:
                    text_parts.append(f"• {page}")
            
            content_placeholder.text = "\n".join(text_parts)
        
        # 기술 스택 슬라이드
        if 'tech_stack' in content and isinstance(content['tech_stack'], list):
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "기술 스택"
            
            text_parts = []
            for tech in content['tech_stack']:
                text_parts.append(f"• {tech}")
            
            content_placeholder.text = "\n".join(text_parts)
        
        # 개발 제안 슬라이드
        if 'development_recommendations' in content:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "개발 제안"
            content_placeholder.text = content['development_recommendations']
        
        # 결론 슬라이드
        if 'conclusion' in content:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            
            title.text = "결론"
            content_placeholder.text = content['conclusion']

    def export_to_pdf(self, content: Dict[str, Any], filename: str = "export", 
                      css: Optional[str] = None, **kwargs) -> str:
        """
        PDF 형식으로 내보내기 (HTML을 PDF로 변환)
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            css (Optional[str]): CSS 스타일시트 경로
            **kwargs: 추가 옵션
            
        Returns:
            str: 생성된 PDF 파일 경로
        """
        # 파일 경로 설정
        output_path = self.output_dir / f"{filename}.pdf"
        
        try:
            # 먼저 HTML로 변환
            html_path = self.export_to_html(content, f"{filename}_temp", css, **kwargs)
            
            # HTML을 PDF로 변환
            HTML(filename=html_path).write_pdf(output_path)
            
            # 임시 HTML 파일 삭제
            os.remove(html_path)
            
            logger.info(f"PDF 파일 생성: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"PDF 내보내기 실패: {str(e)}")
            raise

    def export_to_notion(self, content: Dict[str, Any], parent_id: Optional[str] = None,
                        page_title: Optional[str] = None, api_key: Optional[str] = None) -> str:
        """
        Notion 페이지로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            parent_id (Optional[str]): 노션 부모 페이지 ID (None이면 환경 변수에서 로드)
            page_title (Optional[str]): 노션 페이지 제목 (None이면 콘텐츠에서 추출)
            api_key (Optional[str]): 노션 API 키 (None이면 환경 변수에서 로드)
            
        Returns:
            str: 생성된 노션 페이지 URL
        
        Raises:
            ValueError: 필수 매개변수가 누락된 경우
            Exception: 노션 페이지 생성 실패 시
        """
        try:
            # 필수 매개변수 확인
            if not parent_id:
                parent_id = os.getenv("NOTION_PARENT_PAGE_ID")
                if not parent_id:
                    raise ValueError("부모 페이지 ID가 필요합니다. parent_id 매개변수 또는 NOTION_PARENT_PAGE_ID 환경 변수를 설정하세요.")
            
            # 페이지 제목 설정
            if not page_title:
                website_name = content.get("website", {}).get("name", "웹사이트")
                page_title = f"{website_name} 클론 기획서"
            
            # 노션 클라이언트 생성
            notion_client = NotionClient(api_key)
            
            # 노션 페이지로 내보내기
            logger.info(f"노션 페이지로 내보내기 시작: {page_title}")
            page_id = notion_client.export_website_analysis(content, parent_id, page_title)
            
            # 페이지 URL 생성
            page_url = f"https://www.notion.so/{page_id.replace('-', '')}"
            logger.info(f"노션 페이지로 내보내기 완료: {page_url}")
            
            return page_url
            
        except Exception as e:
            logger.error(f"노션 페이지로 내보내기 실패: {str(e)}")
            raise

    def export_to_google_drive(self, content: Dict[str, Any], filename: str = "export", 
                             folder_name: Optional[str] = None, 
                             make_public: bool = True,
                             export_format: str = 'markdown', 
                             credentials_file: Optional[str] = None,
                             token_file: Optional[str] = None) -> str:
        """
        구글 드라이브로 내보내기
        
        Args:
            content (Dict[str, Any]): 내보낼 콘텐츠
            filename (str): 출력 파일 이름 (확장자 제외)
            folder_name (Optional[str]): 구글 드라이브에 생성할 폴더 이름 (None이면 현재 날짜 사용)
            make_public (bool): 공개 접근 설정 여부
            export_format (str): 드라이브에 업로드할 파일 형식 ('markdown', 'pdf', 'pptx', 'zip')
            credentials_file (Optional[str]): 구글 OAuth 인증 정보 파일 경로
            token_file (Optional[str]): 인증 토큰 저장 파일 경로
            
        Returns:
            str: 구글 드라이브 파일 공유 URL
        
        Raises:
            ValueError: 지원하지 않는 내보내기 형식
            Exception: 내보내기 실패 시
        """
        try:
            # 유효한 내보내기 형식 확인
            valid_formats = ['markdown', 'md', 'pdf', 'pptx', 'zip']
            if export_format.lower() not in valid_formats:
                raise ValueError(f"지원하지 않는 내보내기 형식: {export_format}. 지원 형식: {', '.join(valid_formats)}")
            
            # 지정된 형식으로 파일 내보내기
            if export_format.lower() in ['markdown', 'md']:
                local_file_path = self.export_to_markdown(content, filename)
            elif export_format.lower() == 'pdf':
                local_file_path = self.export_to_pdf(content, filename)
            elif export_format.lower() == 'pptx':
                local_file_path = self.export_to_pptx(content, filename)
            elif export_format.lower() == 'zip':
                local_file_path = self.export_to_zip(content, filename)
            
            # 폴더 이름 생성
            if not folder_name:
                website_name = content.get('website', {}).get('name', '웹사이트')
                folder_name = f"{website_name} 클론 기획서"
            
            # 구글 드라이브 클라이언트 생성
            gdrive_client = GoogleDriveClient(credentials_file, token_file)
            
            # 파일을 구글 드라이브로 내보내기
            logger.info(f"파일 '{local_file_path}'을(를) 구글 드라이브로 내보내는 중...")
            success, file_id, file_url = gdrive_client.export_to_google_drive(
                file_path=local_file_path,
                folder_name=folder_name,
                make_public=make_public
            )
            
            if not success or not file_url:
                raise Exception("구글 드라이브로 내보내기 실패")
            
            logger.info(f"파일을 구글 드라이브로 내보냈습니다. URL: {file_url}")
            return file_url
            
        except Exception as e:
            logger.error(f"구글 드라이브 내보내기 중 오류 발생: {str(e)}")
            raise


# 유틸리티 함수
def export_content(content: Dict[str, Any], formats: List[str], 
                  output_dir: Optional[Union[str, Path]] = None, 
                  filename: str = "export", **kwargs) -> Dict[str, str]:
    """
    여러 형식으로 내용 내보내기를 위한 유틸리티 함수
    
    Args:
        content (Dict[str, Any]): 내보낼 콘텐츠
        formats (List[str]): 내보내기 형식 목록
        output_dir (Optional[Union[str, Path]]): 출력 디렉터리
        filename (str): 기본 파일 이름
        **kwargs: 각 형식별 추가 옵션
        
    Returns:
        Dict[str, str]: 형식별 출력 파일 경로
    """
    manager = ExportManager(output_dir)
    results = {}
    
    for fmt in formats:
        try:
            output_path = manager.export_to_format(
                content=content,
                format_type=fmt,
                filename=filename,
                **kwargs
            )
            results[fmt] = output_path
        except Exception as e:
            logger.error(f"{fmt} 형식으로 내보내기 실패: {str(e)}")
            results[fmt] = None
    
    return results


def markdown_to_html(markdown_path: Union[str, Path], output_path: Optional[Union[str, Path]] = None, 
                    css: Optional[str] = None) -> str:
    """
    마크다운 파일을 HTML로 변환
    
    Args:
        markdown_path (Union[str, Path]): 마크다운 파일 경로
        output_path (Optional[Union[str, Path]]): 출력 HTML 파일 경로 (None이면 자동 생성)
        css (Optional[str]): CSS 스타일시트 경로
        
    Returns:
        str: 생성된 HTML 파일 경로
    """
    markdown_path = Path(markdown_path)
    
    if not output_path:
        output_path = markdown_path.parent / f"{markdown_path.stem}.html"
    else:
        output_path = Path(output_path)
    
    try:
        # 마크다운 파일 읽기
        with open(markdown_path, "r", encoding="utf-8") as f:
            md_content = f.read()
        
        # 마크다운을 HTML로 변환
        html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
        
        # CSS 추가
        if css and os.path.exists(css):
            with open(css, "r", encoding="utf-8") as f:
                css_content = f.read()
            html_template = f"""<!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0">
                <title>{markdown_path.stem}</title>
                <style>
                {css_content}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
        else:
            # 기본 스타일
            html_template = f"""<!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0">
                <title>{markdown_path.stem}</title>
                <style>
                body {{
                    font-family: 'Noto Sans KR', Arial, sans-serif;
                    line-height: 1.6;
                    max-width: 900px;
                    margin: 0 auto;
                    padding: 20px;
                    color: #333;
                }}
                h1, h2, h3, h4, h5, h6 {{
                    margin-top: 1.5em;
                    margin-bottom: 0.5em;
                    color: #1a1a1a;
                }}
                h1 {{ font-size: 2.2em; border-bottom: 1px solid #eee; padding-bottom: 0.3em; }}
                h2 {{ font-size: 1.8em; border-bottom: 1px solid #eee; padding-bottom: 0.3em; }}
                h3 {{ font-size: 1.5em; }}
                h4 {{ font-size: 1.3em; }}
                table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
                th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
                th {{ background-color: #f2f2f2; }}
                img {{ max-width: 100%; height: auto; }}
                code {{ background-color: #f5f5f5; padding: 2px 4px; border-radius: 3px; }}
                pre {{ background-color: #f5f5f5; padding: 1em; overflow-x: auto; border-radius: 3px; }}
                blockquote {{ background-color: #f9f9f9; border-left: 4px solid #ccc; margin: 1em 0; padding: 0.5em 1em; }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
        
        # 파일 작성
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html_template)
        
        logger.info(f"HTML 파일 생성: {output_path}")
        return str(output_path)
        
    except Exception as e:
        logger.error(f"HTML 변환 실패: {str(e)}")
        raise


def html_to_pdf(html_path: Union[str, Path], output_path: Optional[Union[str, Path]] = None) -> str:
    """
    HTML 파일을 PDF로 변환
    
    Args:
        html_path (Union[str, Path]): HTML 파일 경로
        output_path (Optional[Union[str, Path]]): 출력 PDF 파일 경로 (None이면 자동 생성)
        
    Returns:
        str: 생성된 PDF 파일 경로
    """
    html_path = Path(html_path)
    
    if not output_path:
        output_path = html_path.parent / f"{html_path.stem}.pdf"
    else:
        output_path = Path(output_path)
    
    try:
        # HTML을 PDF로 변환
        HTML(filename=str(html_path)).write_pdf(str(output_path))
        
        logger.info(f"PDF 파일 생성: {output_path}")
        return str(output_path)
        
    except Exception as e:
        logger.error(f"PDF 변환 실패: {str(e)}")
        raise


def markdown_to_pdf(markdown_path: Union[str, Path], output_path: Optional[Union[str, Path]] = None,
                   css: Optional[str] = None) -> str:
    """
    마크다운 파일을 PDF로 변환 (마크다운 → HTML → PDF)
    
    Args:
        markdown_path (Union[str, Path]): 마크다운 파일 경로
        output_path (Optional[Union[str, Path]]): 출력 PDF 파일 경로 (None이면 자동 생성)
        css (Optional[str]): CSS 스타일시트 경로
        
    Returns:
        str: 생성된 PDF 파일 경로
    """
    markdown_path = Path(markdown_path)
    
    if not output_path:
        output_path = markdown_path.parent / f"{markdown_path.stem}.pdf"
    else:
        output_path = Path(output_path)
    
    try:
        # 마크다운을 HTML로 변환 (임시 파일)
        with tempfile.NamedTemporaryFile(suffix='.html', delete=False) as temp_html:
            temp_html_path = temp_html.name
        
        html_path = markdown_to_html(markdown_path, temp_html_path, css)
        
        # HTML을 PDF로 변환
        pdf_path = html_to_pdf(html_path, output_path)
        
        # 임시 HTML 파일 삭제
        if os.path.exists(temp_html_path):
            os.remove(temp_html_path)
        
        return pdf_path
        
    except Exception as e:
        logger.error(f"마크다운에서 PDF 변환 실패: {str(e)}")
        raise


def apply_ai_analysis_to_content(content: Dict[str, Any], ai_analysis: Dict[str, Any]) -> Dict[str, Any]:
    """
    AI 분석 결과를 기획서 콘텐츠에 통합
    
    Args:
        content (Dict[str, Any]): 원본 기획서 콘텐츠
        ai_analysis (Dict[str, Any]): AI 분석 결과
        
    Returns:
        Dict[str, Any]: AI 분석이 통합된 기획서 콘텐츠
    """
    updated_content = content.copy()
    
    # AI 분석 섹션 추가
    if 'ai_analysis' not in updated_content:
        updated_content['ai_analysis'] = {}
    
    # 제공된 AI 분석 결과 통합
    for key, value in ai_analysis.items():
        updated_content['ai_analysis'][key] = value
    
    # 개발 제안 업데이트 (AI 분석 내용이 있는 경우)
    if 'recommendations' in ai_analysis:
        if 'development_recommendations' not in updated_content:
            updated_content['development_recommendations'] = ''
        
        ai_recommendations = "\n\n## AI 분석 기반 개선 제안\n\n"
        ai_recommendations += ai_analysis['recommendations']
        
        updated_content['development_recommendations'] += ai_recommendations
    
    # 디자인 분석 업데이트
    if 'design_insights' in ai_analysis and 'design_analysis' in updated_content:
        design = updated_content['design_analysis']
        if isinstance(design, dict):
            # AI 디자인 제안 추가
            design['ai_insights'] = ai_analysis['design_insights']
    
    # 기능 분석 업데이트
    if 'functional_insights' in ai_analysis and 'functional_analysis' in updated_content:
        func = updated_content['functional_analysis']
        if isinstance(func, dict):
            # AI 기능 제안 추가
            func['ai_insights'] = ai_analysis['functional_insights']
    
    logger.info("AI 분석 결과가 기획서 콘텐츠에 통합되었습니다.")
    return updated_content


def generate_ai_analysis_report(content: Dict[str, Any], url: str, 
                               output_dir: Optional[Union[str, Path]] = None,
                               filename: str = "ai_analysis_report") -> str:
    """
    AI 분석 보고서 생성
    
    Args:
        content (Dict[str, Any]): 기획서 콘텐츠
        url (str): 분석할 웹사이트 URL
        output_dir (Optional[Union[str, Path]]): 출력 디렉토리
        filename (str): 출력 파일 이름
        
    Returns:
        str: 생성된 AI 분석 보고서 경로
    """
    manager = ExportManager(output_dir)
    
    # AI 분석용 메타데이터 생성
    ai_meta_path = manager.export_for_ai_analysis(content, url, filename)
    
    # 보고서 경로 설정
    output_path = manager.output_dir / f"{filename}.md"
    
    try:
        # AI 메타데이터 로드
        with open(ai_meta_path, "r", encoding="utf-8") as f:
            ai_metadata = json.load(f)
        
        # 마크다운 보고서 생성
        md_parts = []
        
        md_parts.append(f"# AI 분석 보고서: {ai_metadata.get('title', '웹사이트 분석')}\n")
        md_parts.append(f"URL: {ai_metadata.get('url', url)}\n")
        md_parts.append(f"생성일: {ai_metadata.get('generation_date', '날짜 정보 없음')}\n")
        
        # 웹사이트 기본 정보
        md_parts.append("## 웹사이트 기본 정보\n")
        if 'site_structure' in ai_metadata:
            site = ai_metadata['site_structure']
            md_parts.append(f"- **제목**: {site.get('title', '제목 정보 없음')}")
            md_parts.append(f"- **설명**: {site.get('meta_description', '설명 정보 없음')}")
            
            # 이미지 통계
            if 'images' in site:
                img_stats = site['images']
                md_parts.append(f"- **이미지**: 총 {img_stats.get('count', 0)}개 " + 
                              f"(대체 텍스트 있음: {img_stats.get('with_alt', 0)}개, " +
                              f"없음: {img_stats.get('without_alt', 0)}개)")
            
            # 링크 통계
            if 'links' in site:
                link_stats = site['links']
                md_parts.append(f"- **링크**: 내부 {link_stats.get('internal_count', 0)}개, " +
                              f"외부 {link_stats.get('external_count', 0)}개")
        
        # 분석 제안 영역
        md_parts.append("\n## AI 분석 제안 영역\n")
        if 'ai_analysis_hints' in ai_metadata and 'focus_areas' in ai_metadata['ai_analysis_hints']:
            for area in ai_metadata['ai_analysis_hints']['focus_areas']:
                md_parts.append(f"- {area}")
        
        # 디자인 분석
        if 'design_analysis' in ai_metadata:
            md_parts.append("\n## 디자인 분석\n")
            design = ai_metadata['design_analysis']
            
            if 'color_palette' in design:
                md_parts.append("### 색상 팔레트\n")
                if isinstance(design['color_palette'], list):
                    for color in design['color_palette']:
                        md_parts.append(f"- `{color}`")
                else:
                    md_parts.append(design['color_palette'])
            
            if 'typography' in design:
                md_parts.append("\n### 타이포그래피\n")
                md_parts.append(design['typography'])
            
            if 'layout' in design:
                md_parts.append("\n### 레이아웃\n")
                md_parts.append(design['layout'])
        
        # 페이지 구조
        if 'page_structure' in ai_metadata:
            md_parts.append("\n## 페이지 구조\n")
            for page in ai_metadata['page_structure']:
                if isinstance(page, dict) and 'name' in page:
                    md_parts.append(f"### {page['name']}\n")
                    if 'description' in page:
                        md_parts.append(f"{page['description']}\n")
                    if 'components' in page and isinstance(page['components'], list):
                        md_parts.append("#### 구성 요소\n")
                        for component in page['components']:
                            md_parts.append(f"- {component}")
                else:
                    md_parts.append(f"- {page}\n")
        
        # 개선 제안
        if 'ai_analysis_hints' in ai_metadata and 'expected_improvements' in ai_metadata['ai_analysis_hints']:
            md_parts.append("\n## 개선 제안\n")
            md_parts.append(ai_metadata['ai_analysis_hints']['expected_improvements'])
        
        # 파일 작성
        md_content = "\n".join(md_parts)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md_content)
        
        logger.info(f"AI 분석 보고서 생성: {output_path}")
        return str(output_path)
        
    except Exception as e:
        logger.error(f"AI 분석 보고서 생성 실패: {str(e)}")
        raise 